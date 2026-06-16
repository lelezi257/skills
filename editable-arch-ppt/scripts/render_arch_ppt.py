#!/usr/bin/env python3
"""Render editable architecture PPT slides from a compact JSON spec.

Spec coordinates are inches on a 13.333 x 7.5 canvas. The script intentionally
uses native PowerPoint shapes/text/connectors and injects arrowheads into OOXML
because python-pptx 1.0.x does not expose arrowhead APIs.
"""
from __future__ import annotations

import argparse
import json
import sys
import zipfile
from pathlib import Path
from typing import Any, Dict, Iterable, List, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

PALETTE = {
    "bg": "F6F8FC",
    "infra": "E8EBF0",
    "k8s": "DAEEDF",
    "node": "D0E7D6",
    "pod": "E2EFDA",
    "app": "DDEBF7",
    "platform": "FFF2CC",
    "sandbox": "EBF6FF",
    "runtime": "F4B183",
    "config": "E1EEFA",
    "white": "FFFFFF",
    "text": "262626",
    "muted": "606060",
    "green": "70AD47",
    "blue": "1F5BB5",
    "purple": "7030A0",
    "orange": "C55A11",
    "yellow_border": "BF9000",
    "app_border": "5B9BD5",
    "gray_border": "AAAAAA",
}

FLOW_STYLE = {
    "data": {"color": "purple", "width": 2.5, "dash": False},
    "control": {"color": "blue", "width": 1.35, "dash": True},
    "runtime": {"color": "orange", "width": 1.5, "dash": False},
    "config": {"color": "blue", "width": 1.0, "dash": True},
    "neutral": {"color": "muted", "width": 1.0, "dash": False},
}


def rgb(value: str | Tuple[int, int, int]) -> RGBColor:
    if isinstance(value, tuple):
        return RGBColor(*value)
    value = PALETTE.get(value, value).lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def add_arrowhead(connector, begin: bool = False, end: bool = True) -> None:
    ln = connector._element.xpath(".//a:ln")
    if not ln:
        return
    ln = ln[0]
    for child in list(ln):
        if child.tag in (qn("a:headEnd"), qn("a:tailEnd")):
            ln.remove(child)
    if begin:
        h = OxmlElement("a:headEnd")
        h.set("type", "triangle")
        h.set("w", "med")
        h.set("len", "med")
        ln.append(h)
    if end:
        t = OxmlElement("a:tailEnd")
        t.set("type", "triangle")
        t.set("w", "med")
        t.set("len", "med")
        ln.append(t)


def set_text(shape, text: str, size: float = 10, bold: bool = False, color: str = "text", align: str = "center") -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    for run in p.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_box(slide, item: Dict[str, Any]):
    x, y, w, h = item["x"], item["y"], item["w"], item["h"]
    typ = MSO_SHAPE.ROUNDED_RECTANGLE if item.get("radius") else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(item.get("fill", "white"))
    shp.line.color.rgb = rgb(item.get("outline", "gray_border"))
    shp.line.width = Pt(item.get("outline_width", 1.1))
    if item.get("dash"):
        shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    set_text(
        shp,
        item.get("text", ""),
        size=item.get("font_size", 9),
        bold=item.get("bold", False),
        color=item.get("font_color", "text"),
        align=item.get("align", "center"),
    )
    return shp


def add_textbox(slide, item: Dict[str, Any]):
    box = slide.shapes.add_textbox(Inches(item["x"]), Inches(item["y"]), Inches(item["w"]), Inches(item["h"]))
    set_text(
        box,
        item.get("text", ""),
        size=item.get("font_size", 9),
        bold=item.get("bold", False),
        color=item.get("font_color", "text"),
        align=item.get("align", "center"),
    )
    if fill := item.get("fill"):
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(fill)
        box.fill.transparency = item.get("transparency", 0)
    return box


def add_connector(slide, p1, p2, style: Dict[str, Any], begin: bool = False, end: bool = True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, Inches(p1[0]), Inches(p1[1]), Inches(p2[0]), Inches(p2[1]))
    c.line.color.rgb = rgb(style.get("color", "blue"))
    c.line.width = Pt(style.get("width", 1.3))
    if style.get("dash"):
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    add_arrowhead(c, begin=begin, end=end)
    return c


def add_flow(slide, flow: Dict[str, Any]):
    style = dict(FLOW_STYLE.get(flow.get("kind", "control"), FLOW_STYLE["control"]))
    style.update(flow.get("style", {}))
    points = flow["points"]
    if len(points) < 2:
        raise ValueError(f"flow {flow.get('id')} must have at least two points")
    # Each segment is a native PPT elbow connector. Use few points; prefer corridors over dense segmentation.
    for idx in range(len(points) - 1):
        add_connector(
            slide,
            points[idx],
            points[idx + 1],
            style,
            begin=bool(flow.get("begin_arrow")) and idx == 0,
            end=bool(flow.get("end_arrow", True)) and idx == len(points) - 2,
        )
    if label := flow.get("label"):
        add_textbox(slide, {
            "x": label["x"], "y": label["y"], "w": label["w"], "h": label["h"],
            "text": label["text"], "font_size": label.get("font_size", 7.5),
            "font_color": label.get("font_color", style.get("color", "muted")),
            "fill": label.get("fill", "white"), "transparency": label.get("transparency", 10),
        })


def render(spec: Dict[str, Any], out: Path) -> Dict[str, int]:
    prs = Presentation()
    prs.slide_width = Inches(spec.get("slide_width", 13.333))
    prs.slide_height = Inches(spec.get("slide_height", 7.5))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(spec.get("background", "bg"))

    for item in spec.get("texts", []):
        add_textbox(slide, item)
    for item in spec.get("boxes", []):
        add_box(slide, item)
    for flow in spec.get("flows", []):
        add_flow(slide, flow)

    prs.save(out)
    return validate(out)


def validate(path: Path) -> Dict[str, int]:
    prs = Presentation(path)
    with zipfile.ZipFile(path) as z:
        xml = "\n".join(
            z.read(n).decode(errors="ignore")
            for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
    return {
        "slides": len(prs.slides),
        "objects": sum(len(s.shapes) for s in prs.slides),
        "bentConnector": xml.count("bentConnector"),
        "plain_line": xml.count('prst="line"'),
        "arrowheads": xml.count("<a:tailEnd") + xml.count("<a:headEnd"),
    }


def _interval_overlap(a1: float, a2: float, b1: float, b2: float) -> float:
    lo = max(min(a1, a2), min(b1, b2))
    hi = min(max(a1, a2), max(b1, b2))
    return max(0.0, hi - lo)


def _box_label(box: Dict[str, Any], idx: int) -> str:
    text = str(box.get("text", "")).replace("\n", " / ").strip()
    return text[:48] or f"box#{idx}"


def lint_spec(spec: Dict[str, Any], clearance: float = 0.18) -> Dict[str, Any]:
    """Detect route segments that are likely to cross or visually hug modules.

    The check is intentionally conservative. It assumes the spec's route points
    describe intended corridors and flags axis-aligned segment portions inside a
    component's clearance gutter. Large background/layer boxes are skipped by
    default; set `lint: false` on any decorative container to exclude it.
    """
    boxes: List[Tuple[int, Dict[str, Any]]] = []
    for idx, item in enumerate(spec.get("boxes", []), start=1):
        area = float(item.get("w", 0)) * float(item.get("h", 0))
        if item.get("lint") is False or area > 18:
            continue
        boxes.append((idx, item))

    warnings: List[Dict[str, Any]] = []
    edge_tol = float(spec.get("edge_hug_tolerance", 0.04))
    for flow_idx, flow in enumerate(spec.get("flows", []), start=1):
        points = flow.get("points", [])
        for seg_idx in range(len(points) - 1):
            terminal_segment = seg_idx == 0 or seg_idx == len(points) - 2
            x1, y1 = map(float, points[seg_idx])
            x2, y2 = map(float, points[seg_idx + 1])
            horizontal = abs(y1 - y2) < 1e-6
            vertical = abs(x1 - x2) < 1e-6
            for box_idx, b in boxes:
                rx, ry = float(b["x"]), float(b["y"])
                rw, rh = float(b["w"]), float(b["h"])
                # Edge-hugging is different from endpoint contact. A segment
                # can be outside the box and still unreadable if it lies on the
                # same y/x as a box edge for a visible length.
                edge_issue = None
                edge_overlap = 0.0
                if horizontal and (abs(y1 - ry) <= edge_tol or abs(y1 - (ry + rh)) <= edge_tol):
                    edge_overlap = _interval_overlap(x1, x2, rx, rx + rw)
                    if edge_overlap > clearance * 1.2:
                        edge_issue = "edge_hug_horizontal"
                elif vertical and (abs(x1 - rx) <= edge_tol or abs(x1 - (rx + rw)) <= edge_tol):
                    edge_overlap = _interval_overlap(y1, y2, ry, ry + rh)
                    if edge_overlap > clearance * 1.2:
                        edge_issue = "edge_hug_vertical"
                if edge_issue:
                    warnings.append({
                        "flow": flow.get("id") or f"flow#{flow_idx}",
                        "segment": seg_idx + 1,
                        "kind": flow.get("kind", "control"),
                        "issue": edge_issue,
                        "box": _box_label(b, box_idx),
                        "overlap_inches": round(edge_overlap, 3),
                    })
                    continue
                ex1, ey1 = rx - clearance, ry - clearance
                ex2, ey2 = rx + rw + clearance, ry + rh + clearance
                issue = None
                overlap = 0.0
                if horizontal:
                    if ey1 <= y1 <= ey2:
                        overlap = _interval_overlap(x1, x2, ex1, ex2)
                        if overlap > clearance * 1.2:
                            inside = ry <= y1 <= ry + rh
                            issue = "crosses_box" if inside else "inside_clearance_gutter"
                elif vertical:
                    if ex1 <= x1 <= ex2:
                        overlap = _interval_overlap(y1, y2, ey1, ey2)
                        if overlap > clearance * 1.2:
                            inside = rx <= x1 <= rx + rw
                            issue = "crosses_box" if inside else "inside_clearance_gutter"
                else:
                    # Non-axis-aligned points are unusual for this renderer.
                    sx1, sx2 = sorted((x1, x2))
                    sy1, sy2 = sorted((y1, y2))
                    if _interval_overlap(sx1, sx2, ex1, ex2) and _interval_overlap(sy1, sy2, ey1, ey2):
                        issue = "diagonal_near_box"
                if issue:
                    # Terminal connector segments normally enter/exit their
                    # source/target boxes. Clearance lint focuses on the
                    # intermediate route corridors that should remain visibly
                    # separated from modules.
                    if terminal_segment:
                        continue
                    warnings.append({
                        "flow": flow.get("id") or f"flow#{flow_idx}",
                        "segment": seg_idx + 1,
                        "kind": flow.get("kind", "control"),
                        "issue": issue,
                        "box": _box_label(b, box_idx),
                        "overlap_inches": round(overlap, 3),
                    })
    return {
        "clearance_inches": clearance,
        "checked_boxes": len(boxes),
        "checked_flows": len(spec.get("flows", [])),
        "warnings": warnings,
        "warning_count": len(warnings),
    }


def example_spec() -> Dict[str, Any]:
    return {
        "texts": [
            {"x": 0.25, "y": 0.10, "w": 12.8, "h": 0.35, "text": "Example: Editable Runtime Architecture", "font_size": 18, "bold": True},
            {"x": 0.45, "y": 0.47, "w": 12.35, "h": 0.20, "text": "Purple=data, blue dashed=control, orange=runtime. Native PPT shapes/connectors.", "font_size": 8.8, "font_color": "muted"},
        ],
        "boxes": [
            {"x": 0.25, "y": 0.78, "w": 12.7, "h": 6.35, "text": "Infrastructure", "fill": "infra", "outline": "gray_border", "bold": True, "align": "left", "lint": False},
            {"x": 1.75, "y": 0.98, "w": 9.0, "h": 5.90, "text": "Kubernetes Cluster", "fill": "k8s", "outline": "green", "bold": True, "align": "left", "lint": False},
            {"x": 0.50, "y": 2.78, "w": 1.05, "h": 0.92, "text": "Client", "fill": "white", "outline": "purple", "radius": True, "bold": True},
            {"x": 2.05, "y": 1.30, "w": 3.10, "h": 2.02, "text": "App Layer", "fill": "app", "outline": "app_border", "bold": True, "align": "left", "lint": False},
            {"x": 2.25, "y": 1.75, "w": 1.15, "h": 0.52, "text": "router", "fill": "app", "outline": "app_border", "bold": True},
            {"x": 3.72, "y": 1.75, "w": 1.08, "h": 0.52, "text": "resolver", "fill": "app", "outline": "app_border", "bold": True},
            {"x": 5.55, "y": 1.25, "w": 4.70, "h": 2.17, "text": "Platform Control Plane", "fill": "platform", "outline": "yellow_border", "bold": True, "align": "left", "lint": False},
            {"x": 5.82, "y": 1.70, "w": 1.25, "h": 0.55, "text": "API Server", "fill": "platform", "outline": "yellow_border", "bold": True},
            {"x": 8.88, "y": 1.70, "w": 1.02, "h": 0.55, "text": "State", "fill": "platform", "outline": "yellow_border", "bold": True},
            {"x": 2.10, "y": 3.88, "w": 8.10, "h": 2.55, "text": "Worker Node", "fill": "node", "outline": "green", "bold": True, "align": "left", "lint": False},
            {"x": 5.45, "y": 4.18, "w": 4.28, "h": 1.98, "text": "Worker Pod", "fill": "pod", "outline": "green", "bold": True, "align": "left", "lint": False},
            {"x": 5.72, "y": 4.62, "w": 1.28, "h": 0.44, "text": "runtime helper", "fill": "runtime", "outline": "orange", "bold": True},
            {"x": 6.00, "y": 5.16, "w": 3.34, "h": 0.68, "text": "Sandbox / workload", "fill": "sandbox", "outline": "app_border", "dash": True},
        ],
        "flows": [
            {"kind": "data", "points": [[1.55, 3.24], [1.88, 3.24], [1.88, 2.01], [2.25, 2.01]], "label": {"x": 1.18, "y": 2.36, "w": 1.05, "h": 0.24, "text": "request"}},
            {"kind": "data", "points": [[2.25, 2.15], [1.88, 2.15], [1.88, 3.58], [8.62, 3.58], [8.62, 5.25]], "label": {"x": 4.15, "y": 3.31, "w": 2.75, "h": 0.25, "text": "data path to workload"}},
            {"kind": "control", "points": [[3.40, 1.92], [5.82, 1.92]], "label": {"x": 4.20, "y": 1.55, "w": 1.20, "h": 0.22, "text": "resolve/create"}},
            {"kind": "control", "points": [[7.07, 1.98], [8.88, 1.98]], "begin_arrow": True, "label": {"x": 7.70, "y": 1.38, "w": 1.35, "h": 0.22, "text": "state"}},
            {"kind": "runtime", "points": [[7.00, 4.84], [7.70, 5.25]], "label": {"x": 6.90, "y": 4.33, "w": 1.35, "h": 0.22, "text": "runtime call"}},
        ],
    }


def main(argv: List[str]) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", type=Path, help="JSON spec to render")
    parser.add_argument("--out", type=Path, help="Output .pptx path")
    parser.add_argument("--write-example", type=Path, help="Write an example JSON spec and exit")
    parser.add_argument("--validate", type=Path, help="Validate an existing PPTX and exit")
    parser.add_argument("--lint-spec", type=Path, help="Lint a JSON spec for line clearance/crossing risks and exit")
    parser.add_argument("--clearance", type=float, default=0.18, help="Minimum connector clearance in inches for --lint-spec")
    args = parser.parse_args(argv)

    if args.write_example:
        args.write_example.write_text(json.dumps(example_spec(), ensure_ascii=False, indent=2))
        print(args.write_example)
        return 0
    if args.validate:
        print(json.dumps(validate(args.validate), ensure_ascii=False, indent=2))
        return 0
    if args.lint_spec:
        spec = json.loads(args.lint_spec.read_text())
        print(json.dumps(lint_spec(spec, args.clearance), ensure_ascii=False, indent=2))
        return 0
    if not args.spec or not args.out:
        parser.error("--spec and --out are required unless using --write-example, --validate, or --lint-spec")
    spec = json.loads(args.spec.read_text())
    stats = render(spec, args.out)
    lint = lint_spec(spec)
    print(json.dumps({"out": str(args.out), **stats, "clearance_warnings": lint["warning_count"]}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
